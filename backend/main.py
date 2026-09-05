import json
import os
import re
import shutil
import urllib.parse as urlparse
from typing import Annotated

import docx
import pandas as pd
import pytesseract
import yt_dlp
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from groq import Groq
from PIL import Image
from pptx import Presentation
from pydantic import BaseModel
from pypdf import PdfReader
from youtube_transcript_api import YouTubeTranscriptApi

load_dotenv()

app = FastAPI(title="Lectura AI - Universal Lecture Synthesizer")

# Frontend connection enable
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://lecture-ai-one-black.vercel.app",
        "http://localhost:3000",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_FOLDER = "temp_uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

client = Groq(api_key=os.getenv("GROQ_API_KEY"))


def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def extract_text_from_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text])


def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\n"
    return text


def extract_text_from_excel(path: str) -> str:
    sheets_text = []
    excel_file = pd.ExcelFile(path)
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(path, sheet_name=sheet_name)
        sheets_text.append(f"Sheet: {sheet_name}\n" + df.to_string())
    return "\n\n".join(sheets_text)


def extract_text_from_image(path: str) -> str:
    image = Image.open(path)
    return pytesseract.image_to_string(image)


def transcribe_audio_file(file_path: str) -> str:
    with open(file_path, "rb") as audio:
        transcription = client.audio.transcriptions.create(
            file=(os.path.basename(file_path), audio.read()),
            model="whisper-large-v3",
        )
    return transcription.text



def download_youtube_audio(url: str, output_path: str = "temp_yt.mp3") -> str:
    parsed_url = urlparse.urlparse(url.strip())
    video_id = None

    if parsed_url.hostname in ("youtu.be", "www.youtu.be"):
        video_id = parsed_url.path.lstrip("/")
    elif parsed_url.hostname in ("youtube.com", "www.youtube.com", "m.youtube.com"):
        if parsed_url.path == "/watch":
            video_id = urlparse.parse_qs(parsed_url.query).get("v", [None])[0]
        elif parsed_url.path.startswith("/shorts/"):
            video_id = parsed_url.path.split("/")[2]

    if video_id:
        if "?" in video_id:
            video_id = video_id.split("?")[0]
        if "&" in video_id:
            video_id = video_id.split("&")[0]

        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
            try:
                t = transcript_list.find_transcript(["en", "hi", "en-IN"])
            except Exception:
                t = transcript_list.find_generated_transcript(["en", "hi"])
            data = t.fetch()
            return " ".join([item.get("text", "") for item in data if "text" in item])
        except Exception:
            pass

    ydl_opts = {
        'format': 'bestaudio/best',
        'outtmpl': 'temp_yt.%(ext)s',
        'postprocessors': [{
            'key': 'FFmpegExtractAudio',
            'preferredcodec': 'mp3',
            'preferredquality': '64',
        }],
        'quiet': True,
        'no_warnings': True
    }

    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
            
        actual_path = "temp_yt.mp3"
        if os.path.exists(actual_path):
            text = transcribe_audio_file(actual_path)
            os.remove(actual_path)
            return text
        else:
            raise ValueError("Audio processing failed.")
    except Exception as e:
        raise ValueError(f"Could not extract audio from video: {str(e)}")
    parsed_url = urlparse.urlparse(url.strip())
    video_id = None

    if parsed_url.hostname in ("youtu.be", "www.youtu.be"):
        video_id = parsed_url.path.lstrip("/")
    elif parsed_url.hostname in ("youtube.com", "www.youtube.com", "m.youtube.com"):
        if parsed_url.path == "/watch":
            video_id = urlparse.parse_qs(parsed_url.query).get("v", [None])[0]
        elif parsed_url.path.startswith("/shorts/"):
            video_id = parsed_url.path.split("/")[2]
        elif parsed_url.path.startswith("/embed/"):
            video_id = parsed_url.path.split("/")[2]

    if video_id and "?" in video_id:
        video_id = video_id.split("?")[0]
    if video_id and "&" in video_id:
        video_id = video_id.split("&")[0]

    if not video_id:
        raise ValueError("Invalid YouTube URL. Please provide a valid video link.")

    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        transcript = None

        # 1. Check for manual subtitles (English, Hindi)
        try:
            transcript = transcript_list.find_transcript(["en", "hi", "en-IN"])
        except Exception:
            pass

        # 2. If no manual subtitles, check generated subtitles
        if not transcript:
            try:
                transcript = transcript_list.find_generated_transcript(["en", "hi"])
            except Exception:
                pass

        # 3. Fallback to any transcript available and translate to English
        if not transcript:
            for t in transcript_list:
                transcript = t.translate("en")
                break

        if not transcript:
            raise ValueError("No subtitles or transcript available for this video.")

        data = transcript.fetch()
        return " ".join([item.get("text", "") for item in data if "text" in item])

    except Exception as e:
        print(f"Transcript Error: {e}")
        raise ValueError(f"Unable to process YouTube captions: {e}")


def generate_study_deck(raw_text: str):
    if not raw_text or len(raw_text.strip()) < 5:
        return {
            "status": "error",
            "message": "File me readable text nahi mila. Scanned photo ke bajaye digital text upload karein.",
        }

    truncated_text = raw_text[:12000]

    prompt = f"""
You are an expert academic tutor. Analyze the following content thoroughly and generate comprehensive study notes, flashcards, and a complete quiz.

Input Content:
{truncated_text}

CRITICAL RULES:
1. You MUST generate AT LEAST 5 FLASHCARDS (aim for 5-6).
2. You MUST generate AT LEAST 5 QUIZ QUESTIONS (aim for 5-6). Do NOT generate fewer than 5.
3. Every quiz question must have 4 options and a valid correct_index (0 to 3).

You must return a raw JSON object with this EXACT structure:
{{
  "title": "Comprehensive Topic Title",
  "summary": "A detailed executive summary covering context and core goals.",
  "key_takeaways": [
    "Key takeaway point 1",
    "Key takeaway point 2",
    "Key takeaway point 3",
    "Key takeaway point 4",
    "Key takeaway point 5"
  ],
  "notes_markdown": "# Detailed Lecture Notes\\n\\n## Core Concept 1\\n- In-depth explanation with bullet points\\n\\n## Core Concept 2\\n- Step-by-step mechanisms and examples",
  "flashcards": [
    {{"category": "Core Concept", "front": "Question / Concept 1", "back": "Detailed Answer 1"}},
    {{"category": "Key Detail", "front": "Question / Concept 2", "back": "Detailed Answer 2"}},
    {{"category": "Application", "front": "Question / Concept 3", "back": "Detailed Answer 3"}},
    {{"category": "Analysis", "front": "Question / Concept 4", "back": "Detailed Answer 4"}},
    {{"category": "Summary", "front": "Question / Concept 5", "back": "Detailed Answer 5"}}
  ],
  "quiz": [
    {{
      "question": "Multiple choice question 1?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_index": 0,
      "correct_answer": "Option A",
      "explanation": "Clear explanation for question 1."
    }},
    {{
      "question": "Multiple choice question 2?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_index": 1,
      "correct_answer": "Option B",
      "explanation": "Clear explanation for question 2."
    }},
    {{
      "question": "Multiple choice question 3?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_index": 2,
      "correct_answer": "Option C",
      "explanation": "Clear explanation for question 3."
    }},
    {{
      "question": "Multiple choice question 4?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_index": 3,
      "correct_answer": "Option D",
      "explanation": "Clear explanation for question 4."
    }},
    {{
      "question": "Multiple choice question 5?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_index": 0,
      "correct_answer": "Option A",
      "explanation": "Clear explanation for question 5."
    }}
  ]
}}
"""

    response = client.chat.completions.create(
        model="llama-3.1-8b-versatile",
        messages=[
            {
                "role": "system",
                "content": "You are a professional academic assistant. Respond ONLY with valid JSON conforming to the schema.",
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.2,
        response_format={"type": "json_object"},
    )

    raw_response = response.choices[0].message.content.strip()
    if raw_response.startswith("```"):
        raw_response = re.sub(r"^```[a-zA-Z]*\n?", "", raw_response)
        raw_response = re.sub(r"\n?```$", "", raw_response)

    return json.loads(raw_response)


# ------------- API Endpoints -------------


@app.get("/")
def home():
    return {"message": "Lectura AI Universal API is live! 🚀"}


@app.post("/process-file")
@app.post("/api/process-file")
def process_file(file: Annotated[UploadFile, File()]):
    temp_path = os.path.join(UPLOAD_FOLDER, file.filename)
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    ext = file.filename.split(".")[-1].lower()
    raw_text = ""

    try:
        if ext == "txt":
            with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                raw_text = f.read()
        elif ext == "pdf":
            raw_text = extract_text_from_pdf(temp_path)
        elif ext in ["doc", "docx"]:
            raw_text = extract_text_from_docx(temp_path)
        elif ext in ["ppt", "pptx"]:
            raw_text = extract_text_from_pptx(temp_path)
        elif ext in ["xlsx", "xls", "csv"]:
            raw_text = extract_text_from_excel(temp_path)
        elif ext in ["png", "jpg", "jpeg", "webp"]:
            raw_text = extract_text_from_image(temp_path)
        elif ext in ["mp3", "wav", "m4a", "ogg", "mp4", "mkv", "mov", "webm"]:
            raw_text = transcribe_audio_file(temp_path)
        else:
            raise HTTPException(
                status_code=400, detail=f"File extension '.{ext}' is not supported."
            )

        if not raw_text.strip():
            raise HTTPException(
                status_code=400,
                detail="No readable text or speech found in this file.",
            )

        study_package = generate_study_deck(raw_text)
        return {"status": "success", "data": study_package}

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"FILE ERROR: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


class LinkInput(BaseModel):
    url: str


@app.post("/process-youtube")
@app.post("/api/process-youtube")
def process_youtube(payload: LinkInput):
    try:
        raw_text = download_youtube_audio(payload.url)
        study_package = generate_study_deck(raw_text)
        return {"status": "success", "data": study_package}
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"YOUTUBE ERROR: {e}")
        raise HTTPException(status_code=500, detail=str(e))